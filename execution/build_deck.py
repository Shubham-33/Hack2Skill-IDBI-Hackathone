"""Fill the IDBI Innovate 2026 prototype submission deck — PREMIUM style.

Design language: restrained green + gold palette, white cards floating on soft
shadows with hairline borders, gold accent details, refined typography. Reads
high-end rather than busy. Builds native flow + architecture diagrams, embeds
real prototype screenshots, writes a new .pptx (template untouched).
"""
from __future__ import annotations
import struct
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "Prototype Submission Deck _ IDBI Innovate.pptx"
OUT = ROOT / "Prospect Assist AI - IDBI Innovate Submission Deck.pptx"
SHOTS = ROOT / ".tmp" / "shots"

# premium palette — deep green + gold + ink neutrals
DEEP = RGBColor(0x0A, 0x3D, 0x33); GREEN = RGBColor(0x0A, 0x7D, 0x6B)
GOLD = RGBColor(0xB8, 0x8A, 0x2A); GOLDL = RGBColor(0xC9, 0xA8, 0x5A)
INK = RGBColor(0x24, 0x2B, 0x33); MUTE = RGBColor(0x6B, 0x72, 0x80)
HAIR = RGBColor(0xE0, 0xE6, 0xEA); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINTG = RGBColor(0xEE, 0xF5, 0xF2); TINTGO = RGBColor(0xF8, 0xF1, 0xDF)

prs = Presentation(str(TEMPLATE))
slides = list(prs.slides)


def png_size(p):
    d = open(p, "rb").read(26); return struct.unpack(">II", d[16:24])


def shadow(shape, alpha=18000, blur=70000, dist=28000):
    """Subtle soft drop shadow — the core of the premium feel."""
    spPr = shape._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    spPr.append(parse_xml(
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{blur}" dist="{dist}" dir="5400000" rotWithShape="0">'
        f'<a:srgbClr val="1B2A2A"><a:alpha val="{alpha}"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'))


def _txt(tf, lines, anchor=None):
    if anchor: tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(ln.get("space", 3)); p.space_before = Pt(0)
        if ln.get("level"): p.level = ln["level"]
        r = p.add_run(); r.text = ln["t"]
        r.font.size = Pt(ln.get("sz", 13)); r.font.bold = ln.get("bold", False)
        r.font.italic = ln.get("italic", False)
        r.font.color.rgb = ln.get("color", INK); r.font.name = ln.get("font", "Calibri")


def textbox(slide, l, t, w, h, lines, anchor=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    _txt(tb.text_frame, lines, anchor)
    return tb


def rect(slide, l, t, w, h, fill, line=None, lw=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE, sh=False):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if sh: shadow(sp)
    return sp


def gold_rule(slide, l, t, w=1.0, h=0.045):
    rect(slide, l, t, w, h, GOLD, shape=MSO_SHAPE.RECTANGLE)


def kicker(slide, text):
    """Small letter-spaced gold label + hairline — editorial polish under the title."""
    textbox(slide, 0.42, 1.34, 6.0, 0.28,
            [{"t": text.upper(), "sz": 10.5, "bold": True, "color": GOLD}])
    gold_rule(slide, 0.44, 1.62, 0.9)


def badge(slide, l, t, d, glyph, fill=DEEP, gcolor=GOLDL, sz=13):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background(); sp.shadow.inherit = False
    tf = sp.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = glyph; r.font.size = Pt(sz); r.font.bold = True
    r.font.color.rgb = gcolor; r.font.name = "Calibri"
    return sp


def pcard(slide, l, t, w, h, glyph, title, desc, rail=GREEN):
    """White card, hairline border, soft shadow, left accent rail, gold-on-deep badge."""
    rect(slide, l, t, w, h, WHITE, line=HAIR, lw=1.0, sh=True)
    rect(slide, l, t + 0.12, 0.07, h - 0.24, rail, shape=MSO_SHAPE.RECTANGLE)
    badge(slide, l + 0.24, t + 0.16, 0.44, glyph, DEEP, GOLDL, sz=15)
    textbox(slide, l + 0.82, t + 0.14, w - 0.95, 0.4, [{"t": title, "sz": 13, "bold": True, "color": DEEP}])
    textbox(slide, l + 0.28, t + 0.64, w - 0.5, h - 0.72, [{"t": desc, "sz": 11, "color": INK}])


def panel(slide, l, t, w, h, header, num=None):
    rect(slide, l, t, w, h, WHITE, line=HAIR, lw=1.0, sh=True)
    rect(slide, l, t, w, 0.52, DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, l, t + 0.26, w, 0.26, DEEP, shape=MSO_SHAPE.RECTANGLE)  # square off bottom of header
    if num:
        badge(slide, l + 0.16, t + 0.09, 0.34, num, GOLD, DEEP, sz=13)
        textbox(slide, l + 0.62, t + 0.1, w - 0.7, 0.34, [{"t": header, "sz": 12.5, "bold": True, "color": WHITE}])
    else:
        textbox(slide, l + 0.24, t + 0.1, w - 0.4, 0.34, [{"t": header, "sz": 12.5, "bold": True, "color": WHITE}])


def icon_row(slide, l, t, w, glyph, text, accent=GREEN, sz=11):
    badge(slide, l, t, 0.3, glyph, accent, WHITE, sz=10)
    textbox(slide, l + 0.42, t - 0.04, w - 0.42, 0.6, [{"t": text, "sz": sz, "color": INK}])


def stat(slide, l, t, w, h, number, label, accent=GREEN):
    rect(slide, l, t, w, h, WHITE, line=HAIR, lw=1.0, sh=True)
    big = h >= 0.95
    nsz = 25 if big else 20
    lab_top = t + h - (0.5 if big else 0.32)
    textbox(slide, l, t + 0.06, w, lab_top - t - (0.14 if big else 0.02),
            [{"t": number, "sz": nsz, "bold": True, "color": accent, "align": PP_ALIGN.CENTER}],
            anchor=MSO_ANCHOR.MIDDLE)
    if big:
        gold_rule(slide, l + w / 2 - 0.28, t + h - 0.46, 0.56)
    textbox(slide, l + 0.08, lab_top, w - 0.16, 0.32,
            [{"t": label, "sz": 10, "color": MUTE, "align": PP_ALIGN.CENTER}])


def ribbon(slide, l, t, w, h, text, fill=DEEP, font=WHITE, sz=13, sh=True):
    sp = rect(slide, l, t, w, h, fill, sh=sh)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(sz); r.font.bold = True
    r.font.color.rgb = font; r.font.name = "Calibri"
    return sp


def chip(slide, l, t, text, accent=GREEN):
    """Outlined premium chip — white fill, hairline-accent border, accent text."""
    w = 0.094 * len(text) + 0.36
    sp = rect(slide, l, t, w, 0.36, WHITE, line=accent, lw=1.0)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(10.5); r.font.bold = True
    r.font.color.rgb = accent; r.font.name = "Calibri"
    return l + w


def chip_flow(slide, l, t, items, accent=GREEN, rightmax=9.65, gap=0.14, line_h=0.46):
    x, y = l, t
    for text in items:
        w = 0.094 * len(text) + 0.36
        if x + w > rightmax:
            x = l; y += line_h
        chip(slide, x, y, text, accent)
        x += w + gap
    return y + line_h


def arrow(slide, l, t, w=0.32, color=GOLD):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(0.2))
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit = False


def down(slide, cx, ty):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(cx), Inches(ty), Inches(0.2), Inches(0.18))
    a.fill.solid(); a.fill.fore_color.rgb = GOLD; a.line.fill.background(); a.shadow.inherit = False


def box(slide, l, t, w, h, text, fill=GREEN, font=WHITE, sz=10.5, sh=True):
    sp = rect(slide, l, t, w, h, fill, sh=sh)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(sz); r.font.bold = True
    r.font.color.rgb = font; r.font.name = "Calibri"
    return sp


def caption(slide, l, t, w, text):
    badge(slide, l, t - 0.02, 0.16, "", GOLD, GOLD, sz=1)  # gold dot
    textbox(slide, l + 0.24, t - 0.06, w - 0.24, 0.3, [{"t": text, "sz": 10.5, "bold": True, "color": DEEP}])


def picture(slide, path, l, t, max_w, max_h):
    w, h = png_size(path); ar = w / h
    tw, th = max_w, max_w / ar
    if th > max_h: th, tw = max_h, max_h * ar
    pic = slide.shapes.add_picture(str(path), Inches(l + (max_w - tw) / 2), Inches(t), Inches(tw), Inches(th))
    pic.line.color.rgb = HAIR; pic.line.width = Pt(1.0); shadow(pic, alpha=22000)
    return pic


# ================================================================ Slide 1 — Team
s = slides[0]
for shp in s.shapes:
    if shp.has_text_frame and shp.text_frame.text.strip().startswith("Team Details"):
        vals = {"Team name:": "  Black Bulls",
                "Team leader name:": "  Shubham",
                "Problem Statement:": "  Open Track — AI copilot for channel-sourced retail lending "
                "(conversational AI · product advisory · financial-health scoring · default prediction)"}
        for p in shp.text_frame.paragraphs:
            if p.text.strip() in vals and p.runs:
                r = p.add_run(); r.text = vals[p.text.strip()]
                r.font.size = p.runs[0].font.size or Pt(16); r.font.bold = False; r.font.color.rgb = GREEN

# ================================================================ Slide 2 — Brief
s = slides[1]; kicker(s, "The Idea")
ribbon(s, 0.42, 1.78, 9.16, 0.6,
       "A channel-aware loan-origination & partner-relationship copilot — on top of IDBI's LOS / BRE / LMS", DEEP, sz=13)
textbox(s, 0.44, 2.46, 9.1, 0.3, [{"t": "Retail lending is sourced through DSAs & dealers and runs on TAT + trust. We work all four sides at once:", "sz": 11, "color": MUTE}])
cw, ch = 4.5, 1.08
pcard(s, 0.42, 2.8, cw, ch, "★", "Customer", "Only suitable, clearly-explained, instant offers — no spam, no black box.", GREEN)
pcard(s, 5.08, 2.8, cw, ch, "●", "RM / Bank", "A queue ranked by risk-adjusted profit, auto-scoring & recommendations.", GOLD)
pcard(s, 0.42, 3.96, cw, ch, "◆", "Channel partner (DSA)", "Instant decisions, transparent payouts, tiering — the loyalty moat.", GOLD)
pcard(s, 5.08, 3.96, cw, ch, "▲", "Risk / Regulator", "Explainable, fair, adverse-selection-aware lending with audit trail.", GREEN)
ribbon(s, 0.42, 5.14, 9.16, 0.32, "Lend faster, safer, and win partner loyalty.", GREEN, sz=12.5)

# ================================================================ Slide 3 — Opportunities
s = slides[2]
panel(s, 0.42, 2.1, 4.5, 3.2, "How it's different")
diffs = [("◆", "Channel-first, not customer-only — models the DSA relationship: adverse-selection, TAT, tiering, churn."),
         ("₹", "Bank-grade economics — ranks by P(convert) × (net interest income − expected credit loss)."),
         ("▲", "Explainable — per-decision reason codes + Reg-B / RBI adverse-action codes; fair-lending exclusions."),
         ("●", "Integrates, not replaces — an intelligence layer on top of LOS / BRE / LMS.")]
for i, (g, t) in enumerate(diffs):
    icon_row(s, 0.66, 2.78 + i * 0.62, 4.1, g, t, GREEN, sz=10.5)
panel(s, 5.08, 2.1, 4.5, 3.2, "How it solves the problem")
sol = [("✓", "Right lead → right product → right price, instantly, with a professional offer."),
       ("✓", "Flags risky sourcing before disbursal → fewer NPAs."),
       ("✓", "Transparent, fast, tiered payouts → partner loyalty → more & better volume."),
       ("★", "The flywheel: win the partner → volume → data → sharper models.")]
for i, (g, t) in enumerate(sol):
    icon_row(s, 5.32, 2.78 + i * 0.62, 4.1, g, t, GOLD, sz=10.5)

# ================================================================ Slide 4 — Features
s = slides[3]; kicker(s, "Features")
def feat_panel(l, header, items):
    panel(s, l, 1.78, 4.5, 3.55, header)
    for i, t in enumerate(items):
        icon_row(s, l + 0.24, 2.5 + i * 0.4, 4.1, "✓", t, GREEN, sz=10.6)
feat_panel(0.42, "Customer & RM", [
    "ML conversion-propensity queue (risk-adjusted)",
    "Explainable financial-health score (0–100)",
    "Deterministic eligibility, rate & EMI engine",
    "Suitability-first next-best-product advisory",
    "Editable offer + risk-based pricing in-band",
    "Premium branded PDF + WhatsApp / Gmail",
    "Per-decision reason codes on every score"])
feat_panel(5.08, "Channel & Partner", [
    "Partner/lead default-risk (adverse-selection)",
    "Channel-mix ROI + book concentration (HHI)",
    "Commission + clawback payout engine",
    "Duplicate / fraud lead detection",
    "Instant decision (<50 ms) + SLA / TAT clock",
    "Partner tiering + churn early-warning",
    "Agentic AI assistant (natural language)"])

# ================================================================ Slide 5 — Process flow
s = slides[4]; kicker(s, "Process Flow")
stages = [("1", "DSA submits\nlead"), ("2", "Instant\nindicative decision"), ("3", "ML-ranked\nqueue (RM)"),
          ("4", "Recommend +\neditable offer"), ("5", "Risk-based\nprice + PDF"), ("6", "Disburse +\nlog outcome")]
bw, gap, bh, x0, y = 1.28, 0.22, 1.0, 0.35, 2.35
for i, (num, txt) in enumerate(stages):
    x = x0 + i * (bw + gap)
    fill = GREEN if i in (0, 5) else DEEP
    box(s, x, y, bw, bh, txt, fill=fill, sz=10)
    badge(s, x + bw / 2 - 0.16, y - 0.2, 0.32, num, GOLD, DEEP, sz=12)
    if i < len(stages) - 1:
        arrow(s, x + bw - 0.02, y + bh / 2 - 0.1, w=gap + 0.06)
ribbon(s, x0 + 1.0, y + bh + 0.5, (bw + gap) * 4, 0.5,
       "↻  Outcomes (won / lost / default) retrain the ML models — the compounding data moat", GOLD, DEEP, sz=11)
textbox(s, 0.35, 4.85, 9.3, 0.6, [{"t": "Numbers boundary: the deterministic engine owns every rupee (eligibility, rate, EMI, ECL); ML owns probabilities; the LLM only phrases. RM stays in control — edits re-validate through the engine.", "sz": 10.5, "color": MUTE}])

# ================================================================ Slide 6 — Mockups
s = slides[5]; kicker(s, "Interface")
caption(s, 0.42, 1.82, 4.55, "RM dashboard — risk-adjusted, ML-ranked queue")
picture(s, SHOTS / "dashboard.png", 0.42, 2.05, 4.5, 3.0)
caption(s, 5.08, 1.82, 4.55, "Prospect profile — reason codes & recommendation")
picture(s, SHOTS / "prospect.png", 5.08, 2.05, 4.5, 3.0)

# ================================================================ Slide 7 — Architecture
s = slides[6]; kicker(s, "Architecture")
box(s, 1.6, 1.82, 6.8, 0.46, "Presentation  ·  FastAPI + Jinja2 + Tailwind   (RM dashboard · Partner portal · AI assistant)", fill=DEEP, sz=11)
down(s, 4.9, 2.3)
box(s, 0.9, 2.56, 3.9, 0.46, "Layer 1 · Directives (SOPs in Markdown)", fill=GREEN, sz=11)
box(s, 5.2, 2.56, 3.9, 0.46, "Layer 2 · AI Orchestration — NVIDIA NIM LLM (de-identified, phrasing only)", fill=GREEN, sz=10)
down(s, 4.9, 3.04)
box(s, 0.9, 3.3, 8.2, 0.48, "Layer 3 · Deterministic Engine + ML  —  owns every number", fill=DEEP, sz=12)
subs = ["Eligibility · Rate · EMI", "Economics (ECL / RAROC) + Risk pricing", "ML: Propensity + Default-risk",
        "Explainability (reason codes)", "Commission + Clawback", "PDF (WeasyPrint)"]
sw, sg, sx, sy = 2.63, 0.16, 0.9, 3.86
for i, t in enumerate(subs):
    box(s, sx + (i % 3) * (sw + sg), sy + (i // 3) * 0.46, sw, 0.4, t, fill=WHITE, font=DEEP, sz=9.5, sh=False)
    slist = s.shapes[-1]; slist.line.color.rgb = GREEN; slist.line.width = Pt(1.0)
box(s, 0.9, 4.78, 8.2, 0.38, "Data — products · partners · prospects · history · trained model artifacts (.pkl)", fill=INK, sz=10, sh=False)
textbox(s, 0.9, 5.2, 8.2, 0.3, [{"t": "Integration (PoC): CIBIL bureau · KYC / e-consent (DPDP) · LOS / BRE / LMS · NVIDIA AI Enterprise", "sz": 9.5, "bold": True, "color": MUTE}])

# ================================================================ Slide 8 — Technologies
s = slides[7]; kicker(s, "Technology")
groups = [("Backend & Web", ["Python 3.14", "FastAPI", "Uvicorn", "Jinja2", "Tailwind CSS", "vanilla JS"]),
          ("Machine learning & AI", ["scikit-learn", "NumPy / SciPy", "joblib", "occlusion XAI", "NVIDIA NIM LLM"]),
          ("Documents & PDF", ["WeasyPrint", "pango / cairo"]),
          ("Quality & Deployment", ["pytest · 44 tests", "Playwright", "Docker", "Render / Railway / Cloud Run"])]
y = 1.82
for j, (name, items) in enumerate(groups):
    acc = GREEN if j % 2 == 0 else GOLD
    badge(s, 0.42, y + 0.02, 0.2, "", acc, acc, sz=1)
    textbox(s, 0.7, y - 0.03, 9.0, 0.32, [{"t": name, "sz": 12.5, "bold": True, "color": DEEP}])
    y = chip_flow(s, 0.42, y + 0.4, items, acc) + 0.12

# ================================================================ Slide 9 — Cost
s = slides[8]; kicker(s, "Investment")
rect(s, 0.42, 1.85, 4.5, 2.8, WHITE, line=HAIR, lw=1.0, sh=True)
rect(s, 0.42, 1.85, 4.5, 0.07, GREEN, shape=MSO_SHAPE.RECTANGLE)
textbox(s, 0.42, 2.05, 4.5, 0.36, [{"t": "PROTOTYPE — TODAY", "sz": 11.5, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}])
textbox(s, 0.42, 2.4, 4.5, 0.95, [{"t": "≈ ₹0", "sz": 46, "bold": True, "color": GREEN, "align": PP_ALIGN.CENTER}])
textbox(s, 0.55, 3.5, 4.24, 1.0, [
    {"t": "Open-source stack", "sz": 11, "color": INK, "align": PP_ALIGN.CENTER},
    {"t": "NVIDIA NIM free tier", "sz": 11, "color": INK, "align": PP_ALIGN.CENTER},
    {"t": "Free-tier hosting · no licence cost", "sz": 11, "color": INK, "align": PP_ALIGN.CENTER}])
rect(s, 5.08, 1.85, 4.5, 2.8, WHITE, line=HAIR, lw=1.0, sh=True)
rect(s, 5.08, 1.85, 4.5, 0.07, GOLD, shape=MSO_SHAPE.RECTANGLE)
textbox(s, 5.08, 2.05, 4.5, 0.36, [{"t": "PRODUCTION PoC — INDICATIVE (ANNUAL)", "sz": 11, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}])
for i, t in enumerate(["LLM / GPU inference — usage-based",
                       "Cloud hosting + DB + storage — modest",
                       "Integration — LOS/BRE/KYC/CIBIL (one-time)",
                       "Model-ops — monitoring, retraining, governance"]):
    icon_row(s, 5.32, 2.55 + i * 0.48, 4.1, "▸", t, GREEN, sz=10.6)
ribbon(s, 0.42, 4.82, 9.16, 0.5, "ROI:  cost saved (NPA reduction + payout-leakage → 0 + RM selling hours)  ≫  cost to run", DEEP, sz=12.5)

# ================================================================ Slide 10 — Snapshots
s = slides[9]; kicker(s, "Prototype")
grid = [("offer.png", "Editable offer — live EMI & risk-based pricing"),
        ("channel.png", "Channel intelligence — adverse selection, ROI, HHI"),
        ("submit.png", "Instant decision + adverse-action reason codes"),
        ("deck.png", "Built-in pitch deck (live numbers)")]
xs, ys = [0.42, 5.08], [2.05, 3.85]
for i, (img, cap) in enumerate(grid):
    x, yy = xs[i % 2], ys[i // 2]
    caption(s, x, yy - 0.22, 4.5, cap)
    picture(s, SHOTS / img, x, yy, 4.5, 1.55)

# ================================================================ Slide 11 — Performance
s = slides[10]; kicker(s, "Benchmarks")
cd = CategoryChartData()
cd.categories = ["Conversion\nPropensity", "Partner\nDefault-Risk"]
cd.add_series("AUC ×100", (69, 81))
cd.add_series("Gini", (39, 62))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.42), Inches(1.85), Inches(4.6), Inches(2.35), cd)
ch = gf.chart
ch.has_title = True; ch.chart_title.text_frame.text = "Model performance (hold-out, synthetic — indicative)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout = False
ch.legend.font.size = Pt(9)
try:
    ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb = GREEN
    ch.series[1].format.fill.solid(); ch.series[1].format.fill.fore_color.rgb = GOLD
    for ax in (ch.value_axis, ch.category_axis):
        ax.tick_labels.font.size = Pt(9)
except Exception:
    pass
stat(s, 5.18, 1.85, 2.18, 1.1, "< 50 ms", "Instant decision latency", GREEN)
stat(s, 7.46, 1.85, 2.12, 1.1, "~1–2 s", "LLM reply · 9 s timeout", GOLD)
stat(s, 5.18, 3.08, 4.4, 1.12, "0.22", "Partner-history feature share — no leakage · Bayesian shrinkage", DEEP)
textbox(s, 0.42, 4.34, 9.2, 0.3, [{"t": "Business impact on the synthetic book", "sz": 11.5, "bold": True, "color": DEEP}])
stat(s, 0.42, 4.66, 2.9, 0.7, "₹16.4 Cr", "pipeline value", GREEN)
stat(s, 3.47, 4.66, 2.9, 0.7, "₹0.68 Cr", "risk-adjusted value", GOLD)
stat(s, 6.52, 4.66, 3.06, 0.7, "44 / 44", "automated tests passing", DEEP)

# ================================================================ Slide 12 — Future
s = slides[11]; kicker(s, "Roadmap")
phases = [("Integrate", ["IDBI sandbox APIs", "CIBIL bureau + KYC", "e-consent (DPDP Act)", "LOS / BRE / LMS hooks"]),
          ("Harden", ["Retrain on real conv / NPA / TAT", "SHAP + PSI drift monitoring", "Auth / RBAC + audit trail", "Data residency"]),
          ("Scale", ["Partner mobile portal + WhatsApp", "OTP-gated secure delivery", "Collections & Early-Warning", "NVIDIA AI Enterprise / on-prem"])]
pw = 3.0
for i, (name, items) in enumerate(phases):
    l = 0.42 + i * (pw + 0.18)
    panel(s, l, 1.85, pw, 3.35, f"Phase {i+1} · {name}", num=str(i + 1))
    for j, t in enumerate(items):
        icon_row(s, l + 0.22, 2.6 + j * 0.62, pw - 0.35, "▸", t, GREEN if i != 1 else GOLD, sz=10.4)

# ================================================================ Slide 13 — Links
s = slides[12]
links = [("⎇", "GitHub Public Repository", "https://github.com/Shubham-33/Hack2Skill-IDBI-Hackathone", GREEN),
         ("◉", "Final Product Link (live prototype)", "https://prospect-assist-ai-2guf.onrender.com", GOLD),
         ("►", "Demo Video Link (3 minutes)", "[ paste your unlisted YouTube / Drive link ]", GREEN)]
for i, (g, label, url, rail) in enumerate(links):
    t = 2.15 + i * 1.02
    rect(s, 0.42, t, 9.16, 0.84, WHITE, line=HAIR, lw=1.0, sh=True)
    rect(s, 0.42, t + 0.12, 0.07, 0.6, rail, shape=MSO_SHAPE.RECTANGLE)
    badge(s, 0.66, t + 0.22, 0.42, g, DEEP, GOLDL, sz=15)
    textbox(s, 1.34, t + 0.12, 8.1, 0.35, [{"t": label, "sz": 12.5, "bold": True, "color": DEEP}])
    textbox(s, 1.34, t + 0.46, 8.1, 0.35, [{"t": url, "sz": 12, "color": INK}])

# ================================================================ Slide 15 — Closing
s = slides[14]
phs = [sh for sh in s.shapes if sh.is_placeholder]
if phs:
    phs[0].text_frame.text = "Prospect Assist AI"
    for p in phs[0].text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = GREEN; r.font.bold = True
    if len(phs) > 1:
        phs[1].text_frame.text = "Lend faster, safer, and win partner loyalty.  ·  IDBI Innovate 2026"

prs.save(str(OUT))
print("saved:", OUT.name, "|", len(slides), "slides")
